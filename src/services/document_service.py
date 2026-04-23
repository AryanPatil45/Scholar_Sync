import fitz  # PyMuPDF for fast extraction
import io    # Needed to read PPTX from memory
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

class DocumentService:
    def __init__(self):
        # We limit chunk size to 300 to keep the 3B LLM prompt small and fast on 16GB RAM.
        # We include the Devanagari Danda ("।") to safely tokenize Marathi sentences.
        self.text_splitter = RecursiveCharacterTextSplitter(
            separators=["\n\n", "\n", "।", ".", " ", ""], 
            chunk_size=300,
            chunk_overlap=50,
            length_function=len,
        )

    def process_pdf(self, file_bytes: bytes, filename: str) -> list[dict]:
        """Extracts text from PDF bytes and splits it into manageable chunks."""
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        full_text = ""
        
        for page in doc:
            full_text += page.get_text() + "\n"
            
        doc.close()
        
        chunks = self.text_splitter.split_text(full_text)
        
        processed_chunks = []
        for i, chunk in enumerate(chunks):
            processed_chunks.append({
                "chunk_id": f"{filename}_chunk_{i}",
                "text": chunk,
                "metadata": {"source": filename}
            })
            
        return processed_chunks

    def process_ppt(self, file_bytes: bytes, filename: str) -> list[dict]:
        """Extracts text from PPTX bytes and splits it into manageable chunks."""
        # Wrap the raw bytes in a "file-like" object so python-pptx can read it
        ppt_stream = io.BytesIO(file_bytes)
        presentation = Presentation(ppt_stream)
        
        full_text = ""
        
        # Scrape every single text box on every single slide
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
                    
        # Pass the extracted text to your exact same splitter
        chunks = self.text_splitter.split_text(full_text)
        
        processed_chunks = []
        for i, chunk in enumerate(chunks):
            processed_chunks.append({
                "chunk_id": f"{filename}_chunk_{i}",
                "text": chunk,
                "metadata": {"source": filename}
            })
            
        return processed_chunks